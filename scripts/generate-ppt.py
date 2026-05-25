"""
Generate presentation (7 slides) — 쉬운 언어 버전 v2 (기준 0.85, 달성횟수=전체평균)
Flow:
  1. 연구 질문
  2. 3가지 접근법 로드맵
  3. Branch 1: Flat F1->F4  (한계 발견)
  4. Branch 2: Hierarchy H1->H4  (계층 구조로 극복)
  5. Branch 3: Active A1->A3  (능동 쿼리 개선, A4 제외)
  6. 핵심 비교  A3 = 최종 모델
  7. 결론

Usage: python scripts/generate-ppt.py
Output: out/presentation-v6.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

_BASE   = os.path.normpath(os.path.join(os.path.dirname(__file__), ".."))
SIM_DIR = os.path.join(_BASE, "out", "simulation")

# ── Colours ───────────────────────────────────────────────────────────────────
C_FLAT  = RGBColor(0x6B, 0x9B, 0xD2)   # blue
C_HIER  = RGBColor(0x5C, 0xB8, 0x5C)   # green
C_GOLD  = RGBColor(0xFF, 0xD7, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
C_DGRAY = RGBColor(0x55, 0x55, 0x55)
C_BG    = RGBColor(0x1E, 0x2D, 0x40)   # dark navy
C_RED   = RGBColor(0xE8, 0x53, 0x3A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Primitives ────────────────────────────────────────────────────────────────

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = C_BG
    return sl


def tx(sl, text, l, t, w, h, *,
       size=14, bold=False, italic=False, color=C_WHITE,
       align=PP_ALIGN.LEFT):
    tf = sl.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tf


def bx(sl, l, t, w, h, fill, text=None, *,
       size=13, bold=False, tc=C_WHITE, align=PP_ALIGN.CENTER,
       border=None):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border if border else fill
    if text:
        tf = s.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = tc
    return s


def header(sl, title, subtitle=None, accent=C_GOLD):
    tx(sl, title,
       Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.82),
       size=26, bold=True, color=accent)
    if subtitle:
        tx(sl, subtitle,
           Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.38),
           size=13, italic=True, color=RGBColor(0xBB, 0xCC, 0xDD))
    bx(sl, Inches(0.4), Inches(1.26), Inches(12.5), Emu(4500), accent)


def img(sl, path, l, t, w, h, fallback_label=""):
    ext = os.path.splitext(path)[1].lower()
    if os.path.exists(path) and ext != ".svg":
        sl.shapes.add_picture(path, l, t, w, h)
    else:
        fname = os.path.basename(path)
        b = bx(sl, l, t, w, h, RGBColor(0x1A, 0x28, 0x3C),
               border=RGBColor(0x66, 0x99, 0xCC))
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        note = fallback_label or ""
        r.text = f"[ 그림 삽입 위치 ]\n{fname}\n\n삽입 탭 -> 그림 -> 이 디바이스\nout/simulation/{fname}"
        r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x88, 0xBB, 0xFF)


def table(sl, headers, rows, l, t, w, h, *,
          hfill=C_DGRAY, hl=None, hlfill=C_HIER, fs=12):
    tbl = sl.shapes.add_table(1+len(rows), len(headers), l, t, w, h).table
    for ci, hdr in enumerate(headers):
        c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = hfill
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = C_WHITE
    for ri, row in enumerate(rows):
        hl_row = hl is not None and ri == hl
        bg = hlfill if hl_row else (C_GRAY if ri % 2 == 0 else C_WHITE)
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci); c.fill.solid(); c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs); r.font.bold = hl_row
            r.font.color.rgb = C_WHITE if hl_row else C_BLACK


# ── Slide 1: 연구 질문 ─────────────────────────────────────────────────────────

def s1_question(prs):
    sl = new_slide(prs)
    header(sl, "몇 번 비교하면 사용자 취향을 정확히 파악할 수 있을까?",
           "Slide 1 — 연구 질문")

    # Process flow
    steps = [
        ("사용자가\nA vs B 선택", C_FLAT),
        ("AI가 취향\n패턴 학습",   C_HIER),
        ("취향 가중치\n계산",       C_RED),
        ("정확도\n측정",            RGBColor(0xC0, 0x30, 0x20)),
    ]
    bw, bh = Inches(2.65), Inches(1.1)
    y = Inches(1.6)
    for i, (lbl, c) in enumerate(steps):
        x = Inches(0.4) + i * (bw + Inches(0.38))
        bx(sl, x, y, bw, bh, c, lbl, size=13, bold=True)
        if i < 3:
            tx(sl, "->", x + bw + Emu(28000), y + Inches(0.26),
               Inches(0.37), Inches(0.6),
               size=22, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # Goal box
    bx(sl, Inches(0.4), Inches(2.95), Inches(12.55), Inches(0.55),
       RGBColor(0x2A, 0x42, 0x62),
       "목표: 실제 취향과 AI가 파악한 취향의 일치도 85% 이상 달성에 필요한 평균 비교 횟수를 줄이기",
       size=13, bold=True, tc=C_GOLD)

    bullets = [
        "• 부동산 추천에서 사용자가 매물을 비교하는 횟수는 현실적으로 제한적 (피로도 문제)",
        "• 두 가지 큰 방향(기본 방식 / 계층 방식)으로 나눠 각각 4단계씩 개선 실험",
        "• 각 모델을 20번 반복 시뮬레이션 => 성공률(%) · 평균 달성 횟수 · 최대 정확도 측정",
        "• 참고: Chu & Ghahramani (2005) — 가우시안 프로세스 기반 선호도 학습",
    ]
    for i, b in enumerate(bullets):
        tx(sl, b, Inches(0.4), Inches(3.7) + i * Inches(0.7),
           Inches(12.55), Inches(0.65), size=13, color=C_WHITE)


# ── Slide 2: 3가지 접근법 로드맵 ─────────────────────────────────────────────

def s2_overview(prs):
    sl = new_slide(prs)
    header(sl, "3가지 방향으로 단계적 개선 — 전체 흐름",
           "Slide 2 — 각 방향의 역할")

    branches = [
        ("방향 1\n기본 방식 (Flat)",
         "F1 ~ F4",
         "22가지 항목을 한꺼번에 AI가 학습\nF1: 무작위  F2: 정보량 기반  F3: 초기값  F4: 중요도 비율",
         "한계 발견: 최선도 169.1회 필요",
         C_FLAT, "40%", "169.1회"),
        ("방향 2\n계층 방식",
         "H1 ~ H4",
         "큰 범주 먼저, 세부 항목 나중에 학습\nH1: 기본  H2: 가중치 정규화  H3: 반응 반영  H4: 극단 비교쌍",
         "계층 방식 효과: 27.6회 / 60%",
         C_HIER, "60%", "27.6회"),
        ("방향 3\n질문 방식 개선",
         "A1 ~ A3",
         "H4 기반 + 추가 개선\nA1: 판단 어려운 쌍 선택\nA2: 세부 초기값 개선\nA3: A1+A2 조합 (최종)",
         "성공률 60%, 25.9회",
         C_RED, "60%", "25.9회"),
    ]
    cw = Inches(3.9)
    for i, (title, sub, body, result, c, rate, rnd) in enumerate(branches):
        x = Inches(0.4) + i * (cw + Inches(0.32))
        bx(sl, x, Inches(1.42), cw, Inches(0.65), c, title, size=14, bold=True)
        bx(sl, x, Inches(2.07), cw, Inches(0.38),
           RGBColor(0x2A, 0x3A, 0x50), sub, size=11, tc=C_GOLD)
        bx(sl, x, Inches(2.45), cw, Inches(1.85),
           RGBColor(0x1E, 0x2E, 0x44), body,
           size=11, align=PP_ALIGN.LEFT, tc=C_WHITE)
        bx(sl, x, Inches(4.3), cw, Inches(0.65), c, result, size=12, bold=True)
        bx(sl, x, Inches(4.95), cw*0.5, Inches(0.65),
           RGBColor(0x18, 0x28, 0x3A),
           f"성공률\n{rate}", size=12, bold=True, tc=C_GOLD)
        bx(sl, x+cw*0.5, Inches(4.95), cw*0.5, Inches(0.65),
           RGBColor(0x18, 0x28, 0x3A),
           f"달성 횟수\n{rnd}", size=12, bold=True, tc=C_WHITE)
        if i < 2:
            tx(sl, "->", x+cw+Emu(28000), Inches(3.1),
               Inches(0.3), Inches(0.6),
               size=22, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    bx(sl, Inches(0.4), Inches(5.85), Inches(12.55), Inches(0.55),
       RGBColor(0x1A, 0x38, 0x1A),
       "흐름: 기본 방식의 한계(169.1회) -> 계층 방식으로 27.6회 달성(H4) -> 질문 방식 개선 A3(25.9회)로 최종 적용",
       size=13, bold=True, tc=C_GOLD)

    tx(sl, "참고: Chu&Ghahramani'05 · Duchi'08 · Oh'19 · Sadigh'17 · Rezaei'15",
       Inches(0.4), Inches(7.1), Inches(12.55), Inches(0.32),
       size=9, italic=True, color=C_DGRAY)


# ── Slide 3: Flat 결과 F1→F4 ─────────────────────────────────────────────────

def s3_flat(prs):
    sl = new_slide(prs)
    header(sl, "기본 방식 — F1 -> F4 단계별 개선",
           "Slide 3 — 한계 발견: 22가지 항목을 한꺼번에 학습하는 것의 문제", accent=C_FLAT)

    # Step boxes (top)
    steps = [
        ("F1\n무작위 비교쌍",    "기준선",                           "0%",   "200.0회", "0.725"),
        ("F2\n+ 정보량 기반 선택", "정보를 가장 많이 얻는 쌍 선택\n(Brochu 2007)", "15%", "192.9회", "0.796"),
        ("F3\n+ 초기값 설정",    "AI 초기 상태 개선",                "40%",  "169.1회", "0.833"),
        ("F4\n+ 중요도 비율",    "항목별 중요도를 초기값에 반영\n(Rezaei 2015)",   "5%",  "198.8회", "0.778"),
    ]
    sw = Inches(3.0)
    shade_base = [
        RGBColor(0x3A, 0x5A, 0x8A),
        RGBColor(0x2A, 0x70, 0xB0),
        RGBColor(0x2A, 0x60, 0xA0),
        RGBColor(0x3A, 0x5A, 0x8A),
    ]
    for i, (lbl, desc, rate, rnd, cos) in enumerate(steps):
        x = Inches(0.4) + i * (sw + Inches(0.21))
        bx(sl, x, Inches(1.42), sw, Inches(0.72), shade_base[i], lbl,
           size=13, bold=True)
        tx(sl, desc, x, Inches(2.17), sw, Inches(0.65),
           size=10, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
        bx(sl, x, Inches(2.87), sw*0.5, Inches(0.5),
           RGBColor(0x1E, 0x2E, 0x48),
           f"성공률\n{rate}", size=11, bold=True,
           tc=C_GOLD if rate != "0%" else C_RED)
        bx(sl, x+sw*0.5, Inches(2.87), sw*0.5, Inches(0.5),
           RGBColor(0x1E, 0x2E, 0x48),
           f"달성 횟수\n{rnd}", size=11, bold=True, tc=C_WHITE)
        bx(sl, x, Inches(3.4), sw, Inches(0.38),
           RGBColor(0x18, 0x28, 0x40),
           f"최대 정확도: {cos}", size=10, tc=RGBColor(0xAA, 0xFF, 0xAA))
        if i < 3:
            tx(sl, "->", x+sw+Emu(16000), Inches(2.7),
               Inches(0.19), Inches(0.5),
               size=17, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # Chart PNG
    img(sl, os.path.join(SIM_DIR, "branch-flat.png"),
        Inches(0.4), Inches(3.88), Inches(12.55), Inches(2.55),
        "branch-flat.png")

    # Limitation
    bx(sl, Inches(0.4), Inches(6.52), Inches(12.55), Inches(0.55),
       RGBColor(0x3A, 0x18, 0x18),
       "핵심 한계: 22가지 항목을 한꺼번에 학습 => 정보 부족 문제 => 최선(F3)도 평균 169.1회 필요 (비현실적)",
       size=13, bold=True, tc=C_RED)

    tx(sl, "참고: Chu&Ghahramani (2005) · Brochu et al. (2007) · Rezaei (2015)",
       Inches(0.4), Inches(7.15), Inches(12.55), Inches(0.3),
       size=9, italic=True, color=C_DGRAY)


# ── Slide 4: Hierarchy 결과 H1→H4 ────────────────────────────────────────────

def s4_hier(prs):
    sl = new_slide(prs)
    header(sl, "계층 방식 — H1 -> H4 단계별 개선",
           "Slide 4 — 큰 범주 먼저, 세부 항목 나중에 — 정보 부족 문제 극복", accent=C_HIER)

    # Left: hierarchy diagram (compact)
    bx(sl, Inches(0.4), Inches(1.42), Inches(5.55), Inches(0.55),
       C_HIER, "1단계 — 큰 범주 학습 (4가지 중 무엇을 중시?)", size=12, bold=True)
    cats = ["거리", "가격", "안전", "편의성"]
    for j, c in enumerate(cats):
        bx(sl, Inches(0.4)+j*Inches(1.35), Inches(2.02),
           Inches(1.2), Inches(0.45),
           RGBColor(0x2A, 0x5A, 0x3A), c, size=11, bold=True)
    tx(sl, "-> 가장 중요한 2개 범주 선택",
       Inches(0.4), Inches(2.52), Inches(5.55), Inches(0.36),
       size=11, color=C_GOLD, align=PP_ALIGN.CENTER)
    bx(sl, Inches(0.4), Inches(2.93), Inches(5.55), Inches(0.55),
       C_HIER, "2단계 — 세부 항목 학습 (선택된 범주 안에서)", size=12, bold=True)
    for j, s in enumerate(["도보/버스/경사도", "월세/보증금/관리비"]):
        bx(sl, Inches(0.4)+j*Inches(2.78), Inches(3.53),
           Inches(2.65), Inches(0.42),
           RGBColor(0x1A, 0x3A, 0x22), s, size=10)

    # Right: result table
    hdrs = ["변형", "추가된 점", "성공률", "달성 횟수", "정확도"]
    rows = [
        ["H1", "기본 계층 구조",              "0%",   "52.0회",  "0.703"],
        ["H2", "+ 가중치 합을 1로 유지",       "50%",  "28.8회",  "0.827"],
        ["H3", "+ 이전 답변 반영 조정",        "55%",  "25.8회",  "0.845"],
        ["H4", "+ 극단적 비교쌍 추가  *기반*",  "60%",  "27.6회",  "0.843"],
    ]
    table(sl, hdrs, rows,
          Inches(6.15), Inches(1.42), Inches(6.8), Inches(2.65),
          hfill=C_HIER, hl=3, hlfill=RGBColor(0x1A, 0x7A, 0x1A))

    bx(sl, Inches(6.15), Inches(4.15), Inches(6.8), Inches(0.48),
       RGBColor(0x1A, 0x42, 0x1A),
       "F3: 169.1회  →  H2: 28.8회 (5.9배 감소)  →  H4: 27.6회, 성공률 60%",
       size=12, bold=True, tc=C_GOLD)

    # Virtual item explanation
    bx(sl, Inches(0.4), Inches(4.05), Inches(5.55), Inches(0.45),
       RGBColor(0x2A, 0x3A, 0x55),
       "극단적 비교쌍이란?", size=12, bold=True, tc=C_GOLD)
    tx(sl,
       "실제 매물로 만들기 어려운 극단 조건의 가상 비교쌍을 만들어 활용\n"
       "  예) 가상A: 거리 매우 가까움, 가격 매우 비쌈\n"
       "       가상B: 거리 매우 멈, 가격 매우 쌈\n"
       "  -> 거리 vs 가격 중 어느 것이 더 중요한지 빠르게 파악",
       Inches(0.4), Inches(4.55), Inches(5.55), Inches(0.85),
       size=11, color=C_WHITE)

    # Chart PNG
    img(sl, os.path.join(SIM_DIR, "branch-hierarchy.png"),
        Inches(0.4), Inches(5.45), Inches(12.55), Inches(1.6),
        "branch-hierarchy.png")

    tx(sl, "참고: Duchi et al. (2008) · Oh et al. (2019) · Sadigh et al. (2017) · Furnkranz & Hullermeier (2010)",
       Inches(0.4), Inches(7.15), Inches(12.55), Inches(0.3),
       size=9, italic=True, color=C_DGRAY)


# ── Slide 5: Branch 3 — A1→A3 ───────────────────────────────────────────────

def s5_active(prs):
    sl = new_slide(prs)
    header(sl, "방향 3 — 비교쌍 선택 방식 개선 (A1 -> A3)",
           "Slide 5 — 계층 방식에서 질문 방식을 추가로 개선", accent=C_RED)

    # Step boxes
    steps = [
        ("A1\n판단 어려운 쌍 선택",
         "H4 기반 + 고르기 어려운 쌍을\n우선 질문 (정보 최대화)\n(Sadigh 2017)",
         "50%", "28.6회", "0.839", False),
        ("A2\n+ 세부 초기값 개선",
         "H4 기반 + 항목 중요도를\n세부 학습 초기값에 반영\n(Rezaei 2015)",
         "60%", "24.7회", "0.847", False),
        ("A3\nA1 + A2 조합\n★ 최종 모델",
         "H4 + 판단 어려운 쌍 선택\n+ 세부 초기값 개선\n두 개선 방식을 함께 적용",
         "60%", "25.9회", "0.851", True),
    ]
    sw = Inches(3.9)
    shades = [
        RGBColor(0xB0, 0x40, 0x28),
        RGBColor(0xC8, 0x48, 0x30),
        RGBColor(0xE8, 0x53, 0x3A),
    ]
    for i, (lbl, desc, rate, rnd, cos, is_final) in enumerate(steps):
        x = Inches(0.4) + i * (sw + Inches(0.22))
        header_c = C_GOLD if is_final else shades[i]
        bx(sl, x, Inches(1.42), sw, Inches(0.72), header_c, lbl,
           size=13, bold=True, tc=C_BLACK if is_final else C_WHITE)
        tx(sl, desc, x, Inches(2.17), sw, Inches(0.85),
           size=10, color=RGBColor(0xFF, 0xCC, 0xBB), align=PP_ALIGN.CENTER)
        rate_c = C_GOLD
        bx(sl, x, Inches(3.06), sw*0.5, Inches(0.5),
           RGBColor(0x28, 0x18, 0x18),
           f"성공률\n{rate}", size=11, bold=True, tc=rate_c)
        bx(sl, x+sw*0.5, Inches(3.06), sw*0.5, Inches(0.5),
           RGBColor(0x28, 0x18, 0x18),
           f"달성 횟수\n{rnd}", size=11, bold=is_final,
           tc=C_GOLD if is_final else C_WHITE)
        bx(sl, x, Inches(3.59), sw, Inches(0.38),
           RGBColor(0x20, 0x14, 0x14),
           f"최대 정확도: {cos}", size=10, tc=RGBColor(0xFF, 0xCC, 0xAA))
        if i < 2:
            tx(sl, "->", x+sw+Emu(16000), Inches(2.8),
               Inches(0.2), Inches(0.5),
               size=17, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # Chart PNG (A1-A3, A4 excluded)
    img(sl, os.path.join(SIM_DIR, "branch-active.png"),
        Inches(0.4), Inches(4.1), Inches(12.55), Inches(2.55),
        "branch-active.png")

    # A3 is final model
    bx(sl, Inches(0.4), Inches(6.75), Inches(12.55), Inches(0.45),
       RGBColor(0x38, 0x28, 0x10),
       "A3 = 최종 모델: H4 기반 + 판단 어려운 쌍 선택 + 세부 초기값 개선  =>  성공률 60%, 평균 25.9회",
       size=12, bold=True, tc=C_GOLD)

    tx(sl, "참고: Sadigh et al. (2017) · Rezaei (2015)",
       Inches(0.4), Inches(7.2), Inches(12.55), Inches(0.28),
       size=9, italic=True, color=C_DGRAY)


# ── Slide 6: 핵심 비교 — A3 최종 모델 ────────────────────────────────────────

def s6_compare(prs):
    sl = new_slide(prs)
    header(sl, "최종 비교 — A3가 최적 모델",
           "Slide 6 — 전체 11개 모델 한눈에", accent=C_GOLD)

    # Full chart placeholder (SVG only)
    img(sl, os.path.join(SIM_DIR, "variant-chart.svg"),
        Inches(0.4), Inches(1.42), Inches(12.55), Inches(4.0),
        "variant-chart.svg  (삽입 탭->그림->이 디바이스)")

    # 3 KPI boxes: F3 -> H4 -> A3 (final)
    kpis = [
        ("F3 — 기본 방식 최선",  "성공률  40%\n달성 횟수 169.1회", C_FLAT),
        ("H4 — 계층+극단 비교쌍","성공률  60%\n달성 횟수  27.6회", C_HIER),
        ("A3 — 최종 모델",       "성공률  60%\n달성 횟수  25.9회", C_GOLD),
    ]
    bw = Inches(3.9)
    for i, (ttl, body, c) in enumerate(kpis):
        x = Inches(0.4) + i*(bw + Inches(0.22))
        bx(sl, x, Inches(5.6), bw, Inches(0.44), c, ttl,
           size=12, bold=True, tc=C_BLACK if c == C_GOLD else C_WHITE)
        bx(sl, x, Inches(6.04), bw, Inches(0.76),
           RGBColor(0x22, 0x32, 0x48), body,
           size=14, bold=True, tc=c)
        if i < 2:
            tx(sl, "->", x+bw+Emu(24000), Inches(6.2),
               Inches(0.2), Inches(0.5),
               size=18, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    tx(sl,
       "F3 대비 6.5배 빠른 달성 (169.1회 -> 25.9회)  |  A3: 계층 방식 + 극단 비교쌍 + 판단 어려운 쌍 + 세부 초기값",
       Inches(0.4), Inches(6.9), Inches(12.55), Inches(0.42),
       size=12, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)


# ── Slide 7: 결론 ─────────────────────────────────────────────────────────────

def s7_conclusion(prs):
    sl = new_slide(prs)
    header(sl, "결론 — 계층 방식 + 질문 개선 = 비교 횟수 최소화  (최종: A3)")

    conclusions = [
        (C_FLAT,
         "계층 방식 도입",
         "169.1회 -> 28.8회  (5.9배 감소, H2 기준)",
         "22가지 항목 전체 대신 '큰 범주 먼저, 세부 항목 나중에' 2단계로 나눠 정보 부족 문제 극복"),
        (C_HIER,
         "극단적 비교쌍 추가",
         "28.8회 -> 27.6회  (H4, 성공률 60%)",
         "실제 매물로는 만들기 어려운 극단 조건의 가상 쌍을 추가 -> 범주 간 선호도 빠르게 파악"),
        (C_GOLD,
         "질문 방식 개선  A3",
         "27.6회 -> 25.9회  (최종 모델, 성공률 60%)",
         "F3: 169.1회 · 40%  =>  A3: 25.9회 · 60%  (F3 대비 6.5배 빠름, 성공률도 20%p 향상)"),
    ]
    for i, (c, tag, headline, detail) in enumerate(conclusions):
        y = Inches(1.38) + i * Inches(1.62)
        bx(sl, Inches(0.4), y, Inches(1.85), Inches(1.28), c, tag,
           size=13, bold=True,
           tc=C_BLACK if c == C_GOLD else C_WHITE)
        bx(sl, Inches(2.38), y, Inches(10.57), Inches(0.65),
           RGBColor(0x2A, 0x3A, 0x55), headline,
           size=17, bold=True, tc=c)
        bx(sl, Inches(2.38), y+Inches(0.65), Inches(10.57), Inches(0.63),
           RGBColor(0x1E, 0x2D, 0x42), detail,
           size=12, tc=C_WHITE, align=PP_ALIGN.LEFT)

    # Future work
    bx(sl, Inches(0.4), Inches(6.22), Inches(12.55), Inches(0.44),
       RGBColor(0x28, 0x28, 0x42),
       "향후 과제: 극단 비교쌍을 더 다양하게 만들어 성공률 25% 이상 개선 가능 · 실제 사용자 실험 필요",
       size=12, tc=RGBColor(0xBB, 0xBB, 0xFF))

    tx(sl,
       "참고문헌: Chu&Ghahramani'05 · Brochu'07 · Duchi'08 · Rezaei'15 · "
       "Furnkranz&Hullermeier'10 · Oh'19 · Sadigh'17",
       Inches(0.4), Inches(7.12), Inches(12.55), Inches(0.32),
       size=9, italic=True, color=C_DGRAY)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Generating slides...")
    s1_question(prs);   print("  [1/7] research question")
    s2_overview(prs);   print("  [2/7] 3-branch roadmap")
    s3_flat(prs);       print("  [3/7] flat F1->F4")
    s4_hier(prs);       print("  [4/7] hierarchy H1->H4")
    s5_active(prs);     print("  [5/7] active A1->A3")
    s6_compare(prs);    print("  [6/7] final comparison A3")
    s7_conclusion(prs); print("  [7/7] conclusion")

    out = os.path.normpath(os.path.join(_BASE, "out", "presentation-v7.pptx"))
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"\nSaved: {out}")
    print()
    print("NOTE: Slide 6 needs variant-chart.svg inserted manually.")
    print("  PowerPoint -> Insert tab -> Pictures -> This Device")
    print("  File: out/simulation/variant-chart.svg")


if __name__ == "__main__":
    main()
